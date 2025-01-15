import pandas as pd
import PyPDF2
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import os
import streamlit as st
from neo4j import GraphDatabase
import openai
from io import BytesIO

# Initialize necessary models and variables
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
embedding_model = SentenceTransformer(EMBEDDING_MODEL)
vector_dimension = 384
index = faiss.IndexFlatL2(vector_dimension)  # FAISS index for semantic search
metadata = {}

# Load OpenAI API key securely
openai.api_key = os.getenv("OPENAI_API_KEY")

def load_csv(file_path):
    """Load and preprocess CSV file."""
    return pd.read_csv(file_path)

def extract_text_from_pdf(uploaded_file):
    """Extract text from an uploaded PDF file."""
    text = ""
    pdf_reader = PyPDF2.PdfReader(BytesIO(uploaded_file.read()))
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(uploaded_file):
    """Extract text from PowerPoint presentations."""
    text = ""
    presentation = Presentation(BytesIO(uploaded_file.read()))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def create_embeddings(texts):
    """Generate embeddings for a list of texts."""
    return embedding_model.encode(texts)

def store_in_faiss(texts, embeddings, metadata_dict):
    """Store embeddings and metadata in FAISS."""
    global index, metadata
    index.add(embeddings)
    metadata.update(metadata_dict)

def neo4j_connect(uri, user, password):
    """Connect to a Neo4j database."""
    return GraphDatabase.driver(uri, auth=(user, password))

def insert_graph_data(driver, data):
    """Insert data into Neo4j."""
    with driver.session() as session:
        for node, edges in data.items():
            session.run("MERGE (n:Entity {name: $name})", name=node)
            for edge, target in edges.items():
                session.run(
                    "MATCH (a:Entity {name: $source}), (b:Entity {name: $target}) "
                    "MERGE (a)-[:RELATION {type: $type}]->(b)",
                    source=node, target=target, type=edge
                )

def search_query(query, top_k=5):
    """Perform semantic search on FAISS index."""
    query_embedding = embedding_model.encode([query])
    distances, indices = index.search(query_embedding, top_k)
    return indices[0], distances[0]

def generate_response(query, context):
    """Generate a response using OpenAI GPT based on the user query and retrieved context."""
    prompt = f"""You are a financial assistant. Answer the following query using the provided context.
    
    Context: {context}
    
    Query: {query}
    
    Response:"""
    
    try:
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=prompt,
            max_tokens=300,
            temperature=0.7,
        )
        return response.choices[0].text.strip()
    except openai.error.OpenAIError as e:
        return f"An error occurred: {e}"

def search_and_generate_response(query, top_k=5):
    """Perform semantic search and generate a response using OpenAI GPT."""
    indices, _ = search_query(query, top_k)
    # Retrieve context from FAISS (metadata lookup)
    context = "\n".join([f"Document {i}: {metadata[i]}" for i in indices])
    return generate_response(query, context)

# Streamlit frontend
st.title("Financial Data Chatbot")

uploaded_csv = st.file_uploader("Upload CSV", type=["csv"])
uploaded_pdf = st.file_uploader("Upload PDF", type=["pdf"])
uploaded_ppt = st.file_uploader("Upload PPT", type=["ppt", "pptx"])

if uploaded_csv:
    df = load_csv(uploaded_csv)
    st.write("CSV Loaded:")
    st.dataframe(df.head())
    # Add embeddings for CSV data
    csv_texts = df.apply(lambda row: " ".join(map(str, row)), axis=1).tolist()
    csv_embeddings = create_embeddings(csv_texts)
    store_in_faiss(csv_texts, csv_embeddings, {i: text for i, text in enumerate(csv_texts)})

if uploaded_pdf:
    pdf_text = extract_text_from_pdf(uploaded_pdf)
    st.write("Extracted Text from PDF:")
    st.text_area("PDF Text", pdf_text[:1000], height=200)
    # Add embeddings for PDF text
    pdf_embeddings = create_embeddings([pdf_text])
    store_in_faiss([pdf_text], pdf_embeddings, {len(metadata): pdf_text})

if uploaded_ppt:
    ppt_text = extract_text_from_ppt(uploaded_ppt)
    st.write("Extracted Text from PPT:")
    st.text_area("PPT Text", ppt_text[:1000], height=200)
    # Add embeddings for PPT text
    ppt_embeddings = create_embeddings([ppt_text])
    store_in_faiss([ppt_text], ppt_embeddings, {len(metadata): ppt_text})

query = st.text_input("Enter your query:")
if query:
    response = search_and_generate_response(query)
    st.write("Response:")
    st.write(response)
